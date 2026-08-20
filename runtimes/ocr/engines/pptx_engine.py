"""Native PPTX parser (v1.20.0 Phase 2). python-pptx, not Docling, no OCR."""

from __future__ import annotations

import io
from typing import Optional

from ..documents import MAX_PAGES, DocumentError, decode_payload
from .base import ParsedDocument, ParsedPage, ProgressFn
from .office_common import markdown_table, raise_office_open_error


class PptxEngine:
    name = "pptx"

    def parse(
        self,
        document_base64: str,
        *,
        dpi: Optional[int] = None,
        max_pages: int = MAX_PAGES,
        progress: Optional[ProgressFn] = None,
    ) -> ParsedDocument:
        del dpi
        data = decode_payload(document_base64)
        try:
            from pptx import Presentation
        except Exception as exc:
            raise DocumentError(
                "engine-unavailable",
                f"python-pptx is not installed: {exc}",
            ) from exc

        try:
            presentation = Presentation(io.BytesIO(data))
            slides = list(presentation.slides)
            capped = max(1, min(max_pages, MAX_PAGES, len(slides) or 1))
            chosen = slides[:capped]
            if not chosen:
                pages = (ParsedPage(index=0, text=""),)
                return ParsedDocument(pages=pages, engine=self.name, markdown=None)

            parsed: list[ParsedPage] = []
            for index, slide in enumerate(chosen):
                text = _render_slide(slide)
                parsed.append(ParsedPage(index=index, text=text))
                if progress is not None:
                    progress(index + 1, len(chosen))

            joined = "\n\n".join(page.text for page in parsed).strip()
            return ParsedDocument(
                pages=tuple(parsed),
                engine=self.name,
                markdown=joined or None,
            )
        except DocumentError:
            raise
        except Exception as exc:
            raise_office_open_error("PowerPoint", exc)


def _render_slide(slide: object) -> str:
    parts: list[str] = []
    for shape in getattr(slide, "shapes", []):
        if getattr(shape, "has_table", False):
            table = shape.table
            rows = [
                [cell.text or "" for cell in row.cells]
                for row in table.rows
            ]
            rendered = markdown_table(rows)
            if rendered:
                parts.append(rendered)
            continue
        if getattr(shape, "has_text_frame", False):
            text = (shape.text_frame.text or "").strip()
            if text:
                parts.append(text)
    return "\n\n".join(parts).strip()
